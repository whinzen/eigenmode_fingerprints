from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN

prs = Presentation()
slide_layout = prs.slide_layouts[6]  # blank slide
slide = prs.slides.add_slide(slide_layout)

# Helper to add a box with title + text
def add_box(left, top, width, height, title, text):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame

    # Title
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Inches(0.25)
    p.alignment = PP_ALIGN.CENTER

    # Body
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Inches(0.2)
    p.alignment = PP_ALIGN.CENTER

# Layout positions (5 panels)
w = Inches(2.2)
h = Inches(1.6)
y = Inches(1.5)

# Panels
add_box(Inches(0.3), y, w, h,
        "A. Stimulus",
        "Naturalistic narrative\nWord onsets + sentence boundaries")

add_box(Inches(2.7), y, w, h,
        "B. Embeddings",
        "LLM token embeddings\n(BERT / Qwen)")

add_box(Inches(5.1), y, w, h,
        "C. Transition geometry",
        "Shift\nPrediction error\nCurvature\nSubspace exit")

add_box(Inches(7.5), y, w, h,
        "D. fMRI + eigenmodes",
        "BOLD → cortical eigenmodes\nMode energy across scales")

add_box(Inches(9.9), y, w, h,
        "E. GLM",
        "Regressors → β profiles\nacross eigenmodes")

# Add arrows (simple lines)
from pptx.enum.shapes import MSO_CONNECTOR
for x in [2.5, 4.9, 7.3, 9.7]:
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x), Inches(2.3),
        Inches(x+0.2), Inches(2.3)
    )

# Save
out_path = "workflow_figure.pptx"
prs.save(out_path)

print(f"Saved: {out_path}")